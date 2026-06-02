"""
AI 学术写作路由 - 集成 Agent 的智能写作工作台
支持：提纲生成 / 分段写作 / 润色降重 / Word模板导出
"""

from typing import Optional, List, Dict
from fastapi import APIRouter, Query, HTTPException, Body
from pydantic import BaseModel
import asyncio
import os
import json
from datetime import datetime
from app.services.ai_service import load_ai_config

router = APIRouter()

DATA_DIR = os.path.join(os.path.dirname(__file__), "..", "..", "data")
os.makedirs(DATA_DIR, exist_ok=True)

# ==================== 请求模型 ====================

class GenerateOutlineRequest(BaseModel):
    topic: str
    subject: str = ""
    paper_type: str = "本科毕业论文"  # 课程论文 / 本科毕业论文 / 硕博毕业论文
    word_count: int = 8000
    references: List[Dict] = []

class GenerateSectionRequest(BaseModel):
    topic: str
    outline: List[Dict]
    current_section: Dict
    previous_content: str = ""
    word_count: int = 1500

class PolishTextRequest(BaseModel):
    text: str
    mode: str = "polish"  # polish / academic / shorten / expand / paraphrase
    target_length: str = ""

class GenerateAbstractRequest(BaseModel):
    topic: str
    full_content: str
    lang: str = "zh"  # zh / en

class SearchLiteratureRequest(BaseModel):
    query: str
    limit: int = 15

class ExportWordRequest(BaseModel):
    content: str  # Markdown 内容
    template_id: str  # 模板 ID
    title: str = ""
    author: str = ""
    institution: str = ""


# ==================== 学术论文模板 ====================

ACADEMIC_TEMPLATES = {
    "tsinghua-thesis": {
        "name": "清华大学博士论文",
        "font_name": "宋体",
        "title_size": 22,
        "heading1_size": 16,
        "heading2_size": 14,
        "body_size": 12,
        "line_spacing": 1.5,
        "page_margin_mm": {"top": 30, "bottom": 25, "left": 30, "right": 25},
        "first_line_indent": 2,  # 字符
        "abstract_font": "楷体",
    },
    "csu-thesis": {
        "name": "中南大学学位论文",
        "font_name": "宋体",
        "title_size": 22,
        "heading1_size": 16,
        "heading2_size": 14,
        "body_size": 12,
        "line_spacing": 1.5,
        "page_margin_mm": {"top": 30, "bottom": 25, "left": 30, "right": 25},
        "first_line_indent": 2,
        "abstract_font": "楷体",
    },
    "nature-article": {
        "name": "Nature 期刊文章",
        "font_name": "Times New Roman",
        "title_size": 14,
        "heading1_size": 12,
        "heading2_size": 11,
        "body_size": 10,
        "line_spacing": 2.0,
        "page_margin_mm": {"top": 25, "bottom": 25, "left": 25, "right": 25},
        "first_line_indent": 0,
        "abstract_font": "Times New Roman",
    },
    "ieee-paper": {
        "name": "IEEE 会议论文",
        "font_name": "Times New Roman",
        "title_size": 24,
        "heading1_size": 10,
        "heading2_size": 10,
        "body_size": 10,
        "line_spacing": 1.0,
        "page_margin_mm": {"top": 19, "bottom": 43, "left": 18, "right": 18},
        "first_line_indent": 0,
        "abstract_font": "Times New Roman",
        "columns": 2,
    },
    "undergrad-general": {
        "name": "本科毕业论文通用模板",
        "font_name": "宋体",
        "title_size": 18,
        "heading1_size": 15,
        "heading2_size": 14,
        "body_size": 12,
        "line_spacing": 1.5,
        "page_margin_mm": {"top": 25, "bottom": 20, "left": 30, "right": 25},
        "first_line_indent": 2,
        "abstract_font": "楷体",
    },
}


# ==================== AI 调用辅助 ====================

async def _call_ai(messages: list, temperature: float = 0.7, task_type: str = None) -> str:
    """调用 AI 服务（非流式），集成 Agent 提示词 + 智能路由"""
    from app.services.ai_service import ai_service
    full = ""
    async for chunk in ai_service.chat(messages, temperature=temperature, task_type=task_type):
        full += chunk
    return full


def _build_system_prompt(base: str = "") -> str:
    """构建学术写作专用的系统提示词"""
    return (
        "你是一位资深的学术论文写作专家，具备以下能力：\n"
        "1. 精通中英文学术写作规范\n"
        "2. 熟悉各类学术论文结构（课程论文、学位论文、期刊论文）\n"
        "3. 擅长构建逻辑严密的学术论证\n"
        "4. 能生成符合 GB/T 7714 规范的参考文献\n"
        "5. 语言风格正式、客观、学术化\n"
        + (f"\n当前任务: {base}" if base else "")
    )


# ==================== API 端点 ====================

@router.post("/generate-outline")
async def generate_outline(req: GenerateOutlineRequest):
    """生成论文提纲"""
    ref_text = ""
    if req.references:
        ref_text = "\n".join([
            f"- {r.get('title','')} ({r.get('authors','')}, {r.get('year','')})"
            for r in req.references[:10]
        ])

    prompt = f"""请为以下论文主题生成详细的论文提纲：

**论文主题**：{req.topic}
**论文学科**：{req.subject or '通用'}
**论文类型**：{req.paper_type}
**目标字数**：{req.word_count}字

{f"**参考相关文献**：\n{ref_text}" if ref_text else ""}

请按以下 JSON 格式返回提纲（严格 JSON，不要额外解释）：
{{
  "title": "论文标题",
  "outline": [
    {{"level": 1, "title": "第一章 绪论", "sections": [
      {{"level": 2, "title": "1.1 研究背景", "estimated_words": 800, "description": "..."}},
      {{"level": 2, "title": "1.2 国内外研究现状", "estimated_words": 1500, "description": "..."}}
    ]}},
    {{"level": 1, "title": "第二章 ...", "sections": [...]}}
  ],
  "keywords": ["关键词1", "关键词2", ...],
  "estimated_total_words": 8000
}}

要求：
- {req.paper_type}通常 {req.word_count} 字左右
- 结构完整，层次清晰，符合中文学术论文规范
- 各章节字数分配合理
- 至少包含 5 个一级章节"""

    result = await _call_ai([
        {"role": "system", "content": _build_system_prompt("生成论文提纲")},
        {"role": "user", "content": prompt}
    ], temperature=0.5, task_type="outline")

    try:
        # 尝试提取 JSON
        json_str = result
        if "```json" in result:
            json_str = result.split("```json")[1].split("```")[0]
        elif "```" in result:
            json_str = result.split("```")[1].split("```")[0]
        parsed = json.loads(json_str.strip())
        return {"success": True, "data": parsed}
    except:
        return {"success": True, "data": {"title": req.topic, "outline": [], "keywords": [], "raw": result}}


@router.post("/generate-section")
async def generate_section(req: GenerateSectionRequest):
    """生成/扩写论文某一节"""
    outline_text = "\n".join([
        f"{'#' * (s.get('level',1))} {s['title']}" + 
        f" (约{s.get('estimated_words',1000)}字)"
        for s in req.outline
    ])

    prev = f"\n\n上文已写内容：\n{req.previous_content[:500]}" if req.previous_content else ""

    prompt = f"""请撰写以下论文的「{req.current_section['title']}」章节：

**论文主题**：{req.topic}

**完整提纲**：
{outline_text}

**当前章节描述**：{req.current_section.get('description', '根据提纲撰写')}
**目标字数**：约 {req.word_count} 字{prev}

要求：
1. 学术语言，逻辑严密，论证充分
2. 如有必要，恰当地引用文献
3. 段落之间要有自然的过渡
4. 首次出现专业术语时给出中文全称和英文缩写
5. 直接输出正文内容，不要添加"本章将..."之类的导语"""

    result = await _call_ai([
        {"role": "system", "content": _build_system_prompt("撰写论文章节")},
        {"role": "user", "content": prompt}
    ], temperature=0.6, task_type="section")

    return {"success": True, "content": result}


@router.post("/polish")
async def polish_text(req: PolishTextRequest):
    """文本润色/改写"""
    mode_prompts = {
        "polish": "请润色以下学术文本，保持原意，使语言更流畅、更学术化：",
        "academic": "请将以下文本改写为更正式、更学术的表达，使用专业术语替代口语化表达：",
        "shorten": "请精简以下文本，删除冗余表达，保留核心观点和关键论据：",
        "expand": f"请扩写以下文本，增加更多学术细节和论证，目标长度约{req.target_length or '原长的1.5倍'}：",
        "paraphrase": "请用不同的表达方式改写以下文本，保持原意不变（用于降重）：",
    }

    prompt = mode_prompts.get(req.mode, mode_prompts["polish"])

    result = await _call_ai([
        {"role": "system", "content": _build_system_prompt("学术文本润色")},
        {"role": "user", "content": f"{prompt}\n\n原文：\n{req.text}"}
    ], temperature=0.4, task_type="polish")

    return {"success": True, "content": result}


@router.post("/generate-abstract")
async def generate_abstract(req: GenerateAbstractRequest):
    """生成中英文摘要"""
    lang_hint = "中文摘要" if req.lang == "zh" else "英文摘要 (English abstract)"
    word_hint = "300-500字" if req.lang == "zh" else "200-300 words"

    prompt = f"""请根据以下论文内容，生成{lang_hint}（{word_hint}）：

**论文主题**：{req.topic}

**论文全文**：
{req.full_content[:8000]}

要求：
- 涵盖研究背景、目的、方法、主要结果和结论
- 语言精炼，信息密度高
- 独立成文，不引用正文序号"""

    result = await _call_ai([
        {"role": "system", "content": _build_system_prompt("生成论文摘要")},
        {"role": "user", "content": prompt}
    ], temperature=0.3, task_type="abstract")

    return {"success": True, "content": result}


@router.post("/search-literature")
async def search_literature(req: SearchLiteratureRequest):
    """集成文献检索（Semantic Scholar + CrossRef）"""
    results = []

    # Semantic Scholar
    try:
        import httpx
        async with httpx.AsyncClient(timeout=20) as client:
            resp = await client.get(
                "https://api.semanticscholar.org/graph/v1/paper/search",
                params={
                    "query": req.query,
                    "limit": min(req.limit, 20),
                    "fields": "title,authors,year,abstract,paperId,doi,citationCount,journal"
                }
            )
            if resp.status_code == 200:
                data = resp.json().get("data", [])
                for p in data:
                    authors = [a.get("name", "") for a in (p.get("authors") or [])]
                    results.append({
                        "title": p.get("title", ""),
                        "authors": ", ".join(authors[:5]),
                        "year": p.get("year", ""),
                        "abstract": p.get("abstract", ""),
                        "doi": p.get("doi", ""),
                        "citations": p.get("citationCount", 0),
                        "journal": (p.get("journal") or {}).get("name", ""),
                        "source": "Semantic Scholar",
                        "paperId": p.get("paperId", ""),
                    })
    except Exception as e:
        pass

    # 补充 CrossRef
    try:
        import httpx
        async with httpx.AsyncClient(timeout=20) as client:
            resp = await client.get(
                "https://api.crossref.org/works",
                params={"query": req.query, "rows": min(req.limit, 10)},
                headers={"User-Agent": "AcaSight/2.0"}
            )
            if resp.status_code == 200:
                items = resp.json().get("message", {}).get("items", [])
                existing_titles = {r["title"].lower() for r in results}
                for item in items:
                    title = item.get("title", [""])[0]
                    if title.lower() in existing_titles:
                        continue
                    authors = [
                        f"{a.get('given','')} {a.get('family','')}"
                        for a in item.get("author", [])[:5]
                    ]
                    results.append({
                        "title": title,
                        "authors": ", ".join(authors),
                        "year": item.get("created", {}).get("date-parts", [[0]])[0][0],
                        "abstract": item.get("abstract", ""),
                        "doi": item.get("DOI", ""),
                        "citations": 0,
                        "journal": (item.get("container-title") or [""])[0],
                        "source": "CrossRef",
                    })
    except:
        pass

    # 缓存结果（30min TTL），不立即入库 → 搜索即用
    cache_id = ""
    if results:
        from app.services.literature_service import cache_search_results
        cache_id = cache_search_results(req.query, results)

    return {"success": True, "total": len(results), "results": results, "cache_id": cache_id}


@router.get("/templates")
async def list_templates():
    """列出所有 Word 导出模板"""
    return {
        "templates": [
            {
                "id": tid,
                "name": t["name"],
                "font": t["font_name"],
                "body_size": t["body_size"],
                "line_spacing": t["line_spacing"],
            }
            for tid, t in ACADEMIC_TEMPLATES.items()
        ]
    }


@router.post("/export-word")
async def export_word(req: ExportWordRequest):
    """将 Markdown 内容导出为 Word 文档（套用模板）"""
    template = ACADEMIC_TEMPLATES.get(req.template_id)
    if not template:
        raise HTTPException(400, f"未知模板: {req.template_id}")

    try:
        from docx import Document
        from docx.shared import Pt, Mm, Cm, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.oxml.ns import qn
    except ImportError:
        raise HTTPException(500, "python-docx 未安装，请运行 pip install python-docx")

    doc = Document()

    # 页面设置
    section = doc.sections[0]
    m = template["page_margin_mm"]
    section.top_margin = Mm(m["top"])
    section.bottom_margin = Mm(m["bottom"])
    section.left_margin = Mm(m["left"])
    section.right_margin = Mm(m["right"])

    font_name = template["font_name"]
    body_size = template["body_size"]

    # 设置默认字体
    style = doc.styles["Normal"]
    font = style.font
    font.name = font_name
    font.size = Pt(body_size)
    style.element.rPr.rFonts.set(qn("w:eastAsia"), font_name)
    pf = style.paragraph_format
    pf.line_spacing = template["line_spacing"]
    if template.get("first_line_indent", 0) > 0:
        pf.first_line_indent = Cm(template["first_line_indent"] * 0.35)

    # 解析 Markdown 并写入
    lines = req.content.split("\n")
    i = 0
    while i < len(lines):
        line = lines[i].rstrip()

        if not line:
            doc.add_paragraph("")
            i += 1
            continue

        if line.startswith("# ") and not line.startswith("## "):
            # 一级标题 = 论文标题
            p = doc.add_paragraph()
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run = p.add_run(line[2:])
            run.bold = True
            run.font.size = Pt(template["title_size"])
            run.font.name = font_name
            run._element.rPr.rFonts.set(qn("w:eastAsia"), font_name)
        elif line.startswith("## "):
            p = doc.add_paragraph()
            run = p.add_run(line[3:])
            run.bold = True
            run.font.size = Pt(template["heading1_size"])
            run.font.name = font_name
            run._element.rPr.rFonts.set(qn("w:eastAsia"), font_name)
        elif line.startswith("### "):
            p = doc.add_paragraph()
            run = p.add_run(line[4:])
            run.bold = True
            run.font.size = Pt(template["heading2_size"])
            run.font.name = font_name
            run._element.rPr.rFonts.set(qn("w:eastAsia"), font_name)
        elif line.startswith("**摘要**") or line.startswith("**Abstract**"):
            p = doc.add_paragraph()
            run = p.add_run(line.replace("**", ""))
            run.font.size = Pt(body_size)
            if template.get("abstract_font"):
                run.font.name = template["abstract_font"]
                run._element.rPr.rFonts.set(qn("w:eastAsia"), template["abstract_font"])
        else:
            p = doc.add_paragraph(line)
            p.style = doc.styles["Normal"]

        i += 1

    # 保存文件
    safe_title = req.title.replace("/", "_").replace("\\", "_").strip() or "论文导出"
    filename = f"{safe_title}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
    output_path = os.path.join(DATA_DIR, filename)
    doc.save(output_path)

    return {
        "success": True,
        "filename": filename,
        "path": output_path,
        "download_url": f"/api/pdf/download/{filename}",
        "template": template["name"],
    }


@router.post("/export-word-via-backend")
async def export_word_via_backend(req: ExportWordRequest):
    """将 Markdown 内容导出为 Word 文档（直接返回文件流）"""
    template = ACADEMIC_TEMPLATES.get(req.template_id)
    if not template:
        raise HTTPException(400, f"未知模板: {req.template_id}")

    try:
        from docx import Document
        from docx.shared import Pt, Mm, Cm
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.oxml.ns import qn
    except ImportError:
        raise HTTPException(500, "python-docx 未安装")

    doc = Document()
    section = doc.sections[0]
    m = template["page_margin_mm"]
    section.top_margin = Mm(m["top"])
    section.bottom_margin = Mm(m["bottom"])
    section.left_margin = Mm(m["left"])
    section.right_margin = Mm(m["right"])

    font_name = template["font_name"]
    body_size = template["body_size"]

    style = doc.styles["Normal"]
    style.font.name = font_name
    style.font.size = Pt(body_size)
    style.element.rPr.rFonts.set(qn("w:eastAsia"), font_name)
    style.paragraph_format.line_spacing = template["line_spacing"]
    if template.get("first_line_indent", 0) > 0:
        style.paragraph_format.first_line_indent = Cm(template["first_line_indent"] * 0.35)

    # 简单解析 Markdown
    import re
    lines = req.content.split("\n")
    for line in lines:
        line = line.rstrip()
        if not line:
            doc.add_paragraph("")
            continue
        if line.startswith("# ") and not line.startswith("## "):
            p = doc.add_paragraph()
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run = p.add_run(line[2:])
            run.bold = True
            run.font.size = Pt(template["title_size"])
            run.font.name = font_name
            run._element.rPr.rFonts.set(qn("w:eastAsia"), font_name)
        elif line.startswith("## "):
            p = doc.add_paragraph()
            run = p.add_run(line[3:])
            run.bold = True
            run.font.size = Pt(template["heading1_size"])
            run.font.name = font_name
            run._element.rPr.rFonts.set(qn("w:eastAsia"), font_name)
        elif line.startswith("### "):
            p = doc.add_paragraph()
            run = p.add_run(line[4:])
            run.bold = True
            run.font.size = Pt(template["heading2_size"])
            run.font.name = font_name
            run._element.rPr.rFonts.set(qn("w:eastAsia"), font_name)
        else:
            p = doc.add_paragraph(line)
            p.style = doc.styles["Normal"]

    # 返回文件流
    from fastapi.responses import StreamingResponse
    import io
    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)

    safe_title = req.title.replace("/", "_")[:50] or "论文导出"
    filename = f"{safe_title}.docx"

    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'}
    )


# ==================== 写作工作流 API (方向一 1.2-1.3) ====================

class WritingWorkspaceRequest(BaseModel):
    title: str
    paper_type: str = "本科毕业论文"
    word_count: int = 8000
    material_ids: List[str] = []
    data_mode: str = "knowledge_base"
    reference_paper_ids: List[int] = []


class OutlineReviseRequest(BaseModel):
    outline: List[Dict]
    instructions: str = ""


class SectionWriteRequest(BaseModel):
    topic: str
    outline: List[Dict]
    section_index: int = 0
    current_section: Optional[Dict] = None  # 前端传 {title: '...'} 格式
    previous_content: str = ""
    word_count: int = 1500
    reference_dimensions: List[Dict] = []


class InterruptConfirmRequest(BaseModel):
    section_index: int
    material_type: str = "upload"
    material_path: Optional[str] = None
    chart_config: Optional[Dict] = None


@router.post("/workspace/create")
async def create_workspace(req: WritingWorkspaceRequest):
    """创建写作工作台会话"""
    session_id = f"ws_{datetime.now().strftime('%Y%m%d%H%M%S')}_{os.urandom(4).hex()}"

    material_info = []
    if req.material_ids:
        from app.services.unified_storage_service import get_unified_storage
        svc = get_unified_storage()
        items = svc.list_materials(limit=200)
        for item in items:
            for mid in req.material_ids:
                if mid in item.get("filename", ""):
                    material_info.append(item)

    ref_info = []
    if req.reference_paper_ids:
        from app.database import get_db
        from app.models.paper import Paper
        from sqlalchemy import select
        async for db in get_db():
            for pid in req.reference_paper_ids:
                result = await db.execute(select(Paper).where(Paper.id == pid))
                paper = result.scalar_one_or_none()
                if paper:
                    ref_info.append(paper.to_dict())
            break

    return {
        "session_id": session_id,
        "title": req.title,
        "paper_type": req.paper_type,
        "word_count": req.word_count,
        "data_mode": req.data_mode,
        "materials": material_info,
        "references": ref_info,
        "status": "created",
    }


@router.post("/workspace/{session_id}/outline/stream")
async def stream_outline(session_id: str, req: GenerateOutlineRequest):
    """SSE 流式生成论文大纲（联动 WorkflowEngine 状态机）"""
    from fastapi.responses import StreamingResponse
    from app.services.workflow_engine import get_workflow_engine, WritingFlowStatus

    engine = get_workflow_engine()
    flow = engine.get_writing_flow(session_id)
    if flow:
        engine.transition_writing_flow(session_id, WritingFlowStatus.OUTLINING)

    async def generate():
        ref_text = ""
        if req.references:
            ref_text = "\n".join([
                f"- {r.get('title','')} ({r.get('authors','')}, {r.get('year','')})"
                for r in req.references[:10]
            ])

        prompt = f"""请为以下论文主题生成详细的论文提纲：

**论文主题**：{req.topic}
**论文学科**：{req.subject or '通用'}
**论文类型**：{req.paper_type}
**目标字数**：{req.word_count}字

{f"**参考相关文献**：\n{ref_text}" if ref_text else ""}

请按以下 JSON 格式返回提纲：
{{
  "title": "论文标题",
  "outline": [
    {{"level": 1, "title": "第一章 绪论", "sections": [
      {{"level": 2, "title": "1.1 研究背景", "estimated_words": 800, "description": "..."}}
    ]}}
  ],
  "keywords": ["关键词1", "关键词2"],
  "estimated_total_words": 8000
}}"""

        from app.services.ai_service import ai_service
        full_text = ''
        async for chunk in ai_service.chat(
            [{"role": "system", "content": _build_system_prompt("生成论文提纲")},
             {"role": "user", "content": prompt}],
            temperature=0.5,
            task_type="outline",
        ):
            full_text += chunk
            yield f"data: {json.dumps({'type': 'outline_delta', 'content': chunk}, ensure_ascii=False)}\n\n"

        # Try to parse the full text as JSON and send complete event
        parsed = None
        try:
            json_str = full_text
            if "```json" in full_text:
                json_str = full_text.split("```json")[1].split("```")[0]
            elif "```" in full_text:
                json_str = full_text.split("```")[1].split("```")[0]
            parsed = json.loads(json_str.strip())
        except:
            parsed = {"title": req.topic, "outline": [], "keywords": [], "raw": full_text}

        yield f"data: {json.dumps({'type': 'outline_complete', 'data': parsed}, ensure_ascii=False)}\n\n"

        if flow:
            engine.transition_writing_flow(session_id, WritingFlowStatus.OUTLINE_REVIEW)

    return StreamingResponse(generate(), media_type="text/event-stream")


@router.post("/workspace/{session_id}/section/stream")
async def stream_section(session_id: str, req: SectionWriteRequest):
    """SSE 流式撰写论文章节（支持中断信号，联动 WorkflowEngine）"""
    from fastapi.responses import StreamingResponse
    from app.services.workflow_engine import get_workflow_engine, WritingFlowStatus

    engine = get_workflow_engine()
    flow = engine.get_writing_flow(session_id)
    if flow:
        engine.transition_writing_flow(session_id, WritingFlowStatus.WRITING)

    outline_text = "\n".join([
        f"{'#' * (s.get('level',1))} {s['title']}"
        for s in req.outline
    ])

    # 兼容前端传 current_section 或 section_index
    if req.current_section:
        section_title = req.current_section.get("title", "未知章节")
        section_desc = req.current_section.get("description", "")
        # Try to find section_index from outline by title
        section_index = next((i for i, s in enumerate(req.outline) if s.get("title") == section_title), req.section_index)
    else:
        current = req.outline[req.section_index] if req.section_index < len(req.outline) else {}
        section_title = current.get("title", "未知章节")
        section_desc = current.get("description", "")
        section_index = req.section_index
    section_title = current.get("title", "未知章节")
    section_desc = current.get("description", "")

    dim_text = ""
    if req.reference_dimensions:
        dim_text = "\n\n**相关文献维度数据**：\n"
        for dim in req.reference_dimensions[:5]:
            dim_text += f"- {dim.get('paper_title', '')}: {dim.get('dimension_label', '')} — {dim.get('content', '')[:200]}\n"

    is_data_section = any(kw in section_title for kw in ["实验", "结果", "数据", "图表", "分析", "评价"])

    prev = f"\n\n上文已写内容：\n{req.previous_content[:500]}" if req.previous_content else ""

    prompt = f"""请撰写以下论文的「{section_title}」章节：

**论文主题**：{req.topic}

**完整提纲**：
{outline_text}

**当前章节描述**：{section_desc}
**目标字数**：约 {req.word_count} 字{prev}{dim_text}

要求：
1. 学术语言，逻辑严密，论证充分
2. 如有必要，恰当地引用文献
3. 段落之间要有自然的过渡
4. 直接输出正文内容"""

    async def generate():
        if is_data_section:
            if flow:
                engine.transition_writing_flow(
                    session_id, WritingFlowStatus.INTERRUPTED,
                    interrupt_info={
                        "section_index": section_index,
                        "section_title": section_title,
                        "reason": "该章节涉及数据/插图，请确认素材来源",
                        "options": [
                            {"key": "upload", "label": "自主上传素材"},
                            {"key": "chart", "label": "AI科研绘图"},
                            {"key": "existing", "label": "已有成品图片"},
                        ],
                    },
                )
            yield f"data: {json.dumps({'type': 'interrupt', 'section_index': section_index, 'section_title': section_title, 'message': '该章节涉及数据/插图，请确认素材来源'}, ensure_ascii=False)}\n\n"
            return

        from app.services.ai_service import ai_service
        full_content = ''
        async for chunk in ai_service.chat(
            [{"role": "system", "content": _build_system_prompt("撰写论文章节")},
             {"role": "user", "content": prompt}],
            temperature=0.6,
            task_type="section",
        ):
            full_content += chunk
            yield f"data: {json.dumps({'type': 'section_delta', 'content': chunk}, ensure_ascii=False)}\n\n"

        yield f"data: {json.dumps({'type': 'section_complete', 'content': full_content, 'section_index': section_index}, ensure_ascii=False)}\n\n"

        if flow:
            is_last = section_index >= len(req.outline) - 1
            if is_last:
                engine.transition_writing_flow(session_id, WritingFlowStatus.COMPLETED)

    return StreamingResponse(generate(), media_type="text/event-stream")


@router.post("/workspace/{session_id}/interrupt/confirm")
async def confirm_interrupt(session_id: str, req: InterruptConfirmRequest):
    """确认中断：用户选择素材后恢复写作（联动 WorkflowEngine）"""
    from app.services.workflow_engine import get_workflow_engine, WritingFlowStatus

    engine = get_workflow_engine()
    flow = engine.get_writing_flow(session_id)
    if flow:
        engine.transition_writing_flow(session_id, WritingFlowStatus.CONFIRMED)
    material_info = {}
    if req.material_type == "upload" and req.material_path:
        material_info = {"type": "upload", "path": req.material_path}
    elif req.material_type == "chart" and req.chart_config:
        material_info = {"type": "chart", "config": req.chart_config}
    elif req.material_type == "existing":
        material_info = {"type": "existing", "path": req.material_path}

    return {
        "session_id": session_id,
        "section_index": req.section_index,
        "confirmed": True,
        "material": material_info,
    }


# ==================== 研究方向 & 试验方案 API (方向一 1.2扩展) ====================

class ResearchDirectionRequest(BaseModel):
    topic: str
    subject: str = ""
    background: str = ""
    existing_literature: List[Dict] = []
    count: int = 5


class ExperimentDesignRequest(BaseModel):
    topic: str
    research_question: str = ""
    methodology: str = ""
    variables: List[str] = []
    constraints: str = ""
    existing_data: str = ""


@router.post("/research-direction")
async def generate_research_directions(req: ResearchDirectionRequest):
    """基于主题和已有文献，生成可行的研究方向"""
    lit_text = ""
    if req.existing_literature:
        lit_text = "\n\n**已有相关文献**：\n" + "\n".join([
            f"- {r.get('title', '')} ({r.get('authors', '')}, {r.get('year', '')})"
            for r in req.existing_literature[:10]
        ])

    prompt = f"""请基于以下信息，生成 {req.count} 个可行的研究方向：

**研究主题**：{req.topic}
**学科领域**：{req.subject or '通用'}
**研究背景**：{req.background or '暂无'}
{lit_text}

请按以下 JSON 格式返回（严格 JSON，不要额外解释）：
{{
  "directions": [
    {{
      "title": "研究方向标题",
      "description": "简要描述该方向的核心问题和研究价值（100-200字）",
      "novelty": "创新点说明",
      "feasibility": "可行性分析",
      "key_questions": ["关键研究问题1", "关键研究问题2"],
      "suggested_methods": ["建议方法1", "建议方法2"],
      "difficulty": "中等/较高/高",
      "related_fields": ["相关交叉领域1", "相关交叉领域2"]
    }}
  ]
}}

要求：
1. 每个方向应有明确的创新性和学术价值
2. 方向之间应有差异化，避免重叠
3. 考虑可行性，不要过于空泛
4. 结合已有文献的空白和不足提出方向"""

    result = await _call_ai([
        {"role": "system", "content": _build_system_prompt("生成研究方向")},
        {"role": "user", "content": prompt}
    ], temperature=0.7, task_type="research_direction")

    try:
        json_str = result
        if "```json" in result:
            json_str = result.split("```json")[1].split("```")[0]
        elif "```" in result:
            json_str = result.split("```")[1].split("```")[0]
        parsed = json.loads(json_str.strip())
        return {"success": True, "data": parsed}
    except:
        return {"success": True, "data": {"directions": [], "raw": result}}


@router.post("/experiment-design")
async def generate_experiment_design(req: ExperimentDesignRequest):
    """基于研究问题生成试验/实验方案"""
    vars_text = ""
    if req.variables:
        vars_text = "\n**已识别的关键变量**：" + "、".join(req.variables)

    data_text = ""
    if req.existing_data:
        data_text = f"\n**已有数据/素材**：{req.existing_data}"

    prompt = f"""请为以下研究设计详细的实验/试验方案：

**研究主题**：{req.topic}
**核心研究问题**：{req.research_question or '待明确'}
**研究方法倾向**：{req.methodology or '不限'}{vars_text}{data_text}
**约束条件**：{req.constraints or '无特殊约束'}

请按以下 JSON 格式返回（严格 JSON，不要额外解释）：
{{
  "experiment_design": {{
    "title": "实验方案标题",
    "objective": "实验目的",
    "hypothesis": "研究假设",
    "design_type": "实验设计类型（如：对照实验、准实验、仿真实验、案例分析等）",
    "variables": {{
      "independent": ["自变量1", "自变量2"],
      "dependent": ["因变量1", "因变量2"],
      "controlled": ["控制变量1", "控制变量2"]
    }},
    "procedure": [
      {{"step": 1, "title": "步骤标题", "description": "详细描述", "duration": "预计耗时"}},
      {{"step": 2, "title": "...", "description": "...", "duration": "..."}}
    ],
    "data_collection": {{
      "methods": ["数据采集方法1", "数据采集方法2"],
      "instruments": ["测量工具/仪器1", "测量工具/仪器2"],
      "sample_size": "建议样本量及理由"
    }},
    "analysis_plan": {{
      "methods": ["分析方法1（如ANOVA、回归分析等）"],
      "tools": ["分析工具1（如SPSS、Python等）"],
      "significance_level": "显著性水平"
    }},
    "validity": {{
      "internal": "内部效度保障措施",
      "external": "外部效度/推广性说明"
    }},
    "ethics": "伦理考量（如涉及）",
    "timeline": "预计时间安排",
    "risks": ["潜在风险1", "潜在风险2"],
    "alternatives": ["备选方案1", "备选方案2"]
  }}
}}

要求：
1. 方案应具体可执行，不要过于笼统
2. 步骤描述应包含足够的操作细节
3. 数据采集和分析方法应与研究问题匹配
4. 考虑可行性和伦理问题"""

    result = await _call_ai([
        {"role": "system", "content": _build_system_prompt("设计实验方案")},
        {"role": "user", "content": prompt}
    ], temperature=0.5, task_type="experiment_design")

    try:
        json_str = result
        if "```json" in result:
            json_str = result.split("```json")[1].split("```")[0]
        elif "```" in result:
            json_str = result.split("```")[1].split("```")[0]
        parsed = json.loads(json_str.strip())
        return {"success": True, "data": parsed}
    except:
        return {"success": True, "data": {"experiment_design": {}, "raw": result}}


# ==================== PPT 生成 API (方向六 6.6) ====================

class PptGenerateRequest(BaseModel):
    title: str
    subject: str = ""
    outline: Optional[List[Dict]] = None
    content: Optional[str] = None
    style: str = "academic"
    slide_count: int = 15


@router.post("/generate-ppt")
async def generate_ppt(req: PptGenerateRequest):
    """基于论文内容生成学术演示文稿"""
    if req.outline:
        sections_text = "\n".join([
            f"- {s.get('title', '')}: {s.get('description', '')}"
            for s in req.outline
        ])
    elif req.content:
        sections_text = req.content[:3000]
    else:
        sections_text = ""

    prompt = f"""请为以下学术主题生成演示文稿(PPT)的详细内容：

**主题**：{req.title}
**学科**：{req.subject or '通用'}
**风格**：{req.style}
**目标页数**：{req.slide_count}页
{f'**大纲**：{sections_text}' if sections_text else ''}

请按以下 JSON 格式返回（严格 JSON，不要额外解释）：
{{
  "slides": [
    {{
      "type": "title|content|two_column|image|table|conclusion",
      "title": "幻灯片标题",
      "content": ["要点1", "要点2", "要点3"],
      "left_column": ["左栏内容"],
      "right_column": ["右栏内容"],
      "notes": "演讲者备注"
    }}
  ],
  "theme": {{
    "primary_color": "#1a5276",
    "secondary_color": "#2980b9",
    "font_title": "Microsoft YaHei",
    "font_body": "Microsoft YaHei"
  }}
}}

要求：
1. 第一张必须是标题页（type=title），包含主题和副标题
2. 第二张是目录/大纲页
3. 中间页覆盖：研究背景、目的、方法、结果、讨论
4. 倒数第二张是结论/总结
5. 最后一张是致谢/Q&A
6. 每页内容精炼，适合演示（不要大段文字）
7. 总页数约{req.slide_count}页"""

    result = await _call_ai([
        {"role": "system", "content": _build_system_prompt("生成学术PPT内容")},
        {"role": "user", "content": prompt}
    ], temperature=0.5, task_type="agent_reasoning")

    try:
        json_str = result
        if "```json" in result:
            json_str = result.split("```json")[1].split("```")[0]
        elif "```" in result:
            json_str = result.split("```")[1].split("```")[0]
        parsed = json.loads(json_str.strip())
        slides = parsed.get("slides", [])
        theme = parsed.get("theme", {})

        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN

            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            primary = RGBColor.from_string((theme.get("primary_color", "#1a5276")).lstrip("#"))
            secondary = RGBColor.from_string((theme.get("secondary_color", "#2980b9")).lstrip("#"))

            for slide_data in slides:
                slide_type = slide_data.get("type", "content")
                slide_layout = prs.slide_layouts[6] if slide_type == "title" else prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)

                title_shape = slide.shapes.title
                if title_shape:
                    title_shape.text = slide_data.get("title", "")
                    for para in title_shape.text_frame.paragraphs:
                        for run in para.runs:
                            run.font.size = Pt(32)
                            run.font.color.rgb = primary
                            run.font.bold = True

                if slide_type == "content":
                    body = slide_data.get("content", [])
                    if hasattr(slide, 'placeholders') and len(slide.placeholders) > 1:
                        body_shape = slide.placeholders[1]
                        tf = body_shape.text_frame
                        tf.clear()
                        for i, item in enumerate(body):
                            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                            p.text = str(item)
                            p.font.size = Pt(18)
                            p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                            p.space_after = Pt(8)
                    else:
                        left = Inches(1)
                        top = Inches(2)
                        width = Inches(11)
                        height = Inches(5)
                        txBox = slide.shapes.add_textbox(left, top, width, height)
                        tf = txBox.text_frame
                        for i, item in enumerate(body):
                            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                            p.text = f"• {item}"
                            p.font.size = Pt(18)
                            p.space_after = Pt(8)

                elif slide_type == "two_column":
                    left_items = slide_data.get("left_column", [])
                    right_items = slide_data.get("right_column", [])
                    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(5.5), Inches(5))
                    right_box = slide.shapes.add_textbox(Inches(7), Inches(2), Inches(5.5), Inches(5))
                    for i, item in enumerate(left_items):
                        p = left_box.text_frame.paragraphs[0] if i == 0 else left_box.text_frame.add_paragraph()
                        p.text = f"• {item}"
                        p.font.size = Pt(16)
                    for i, item in enumerate(right_items):
                        p = right_box.text_frame.paragraphs[0] if i == 0 else right_box.text_frame.add_paragraph()
                        p.text = f"• {item}"
                        p.font.size = Pt(16)

                notes = slide_data.get("notes", "")
                if notes:
                    slide.notes_slide.notes_text_frame.text = notes

            output_dir = os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(os.path.dirname(__file__)))), "exports")
            os.makedirs(output_dir, exist_ok=True)
            safe_title = "".join(c for c in req.title if c.isalnum() or c in " _-")[:50]
            filename = f"{safe_title}.pptx"
            output_path = os.path.join(output_dir, filename)
            prs.save(output_path)

            return {
                "success": True,
                "data": {
                    "filename": filename,
                    "path": output_path,
                    "slide_count": len(slides),
                    "theme": theme,
                },
            }
        except ImportError:
            return {
                "success": True,
                "data": {
                    "slides": slides,
                    "theme": theme,
                    "slide_count": len(slides),
                    "note": "python-pptx 未安装，返回JSON内容。安装后可生成.pptx文件。",
                },
            }

    except Exception as e:
        return {"success": True, "data": {"raw": result, "parse_error": str(e)}}


@router.get("/download-ppt")
async def download_ppt(path: str = Query(..., description="PPT文件路径")):
    """下载生成的PPT文件"""
    import os
    from fastapi.responses import FileResponse
    if not os.path.exists(path):
        raise HTTPException(status_code=404, detail=f"文件不存在: {path}")
    filename = os.path.basename(path)
    return FileResponse(path, filename=filename, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")
